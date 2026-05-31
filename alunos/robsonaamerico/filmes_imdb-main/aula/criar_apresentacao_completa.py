# criar_apresentacao_completa.py
# Instale as dependências: pip install python-pptx

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
from lxml import etree

# ============================================================
# CONTEÚDO DOS SCRIPTS
# ============================================================

scripts_data = [
    {
        "titulo_slide": "Script: Exploração Inicial dos Dados",
        "subtitulo": "script__explorao_.py — Exercício 1",
        "codigo": '''# EXERCÍCIO 1 — Exploração Inicial

import pandas as pd

caminho = r"filmes_imdb\\dados\\dados_excel\\filmes_imdb.xlsx"
df_imdb = pd.read_excel(caminho)

# 1.1 Visualizando as primeiras linhas
print("=== PRIMEIRAS 10 LINHAS ===")
print(df_imdb.head(10))

# 1.2 Verificando informações estruturais
print("\\n=== INFORMAÇÕES DO DATAFRAME ===")
print(df_imdb.info())

# 1.3 Verificando o tamanho da tabela
print("\\n=== TAMANHO DA TABELA (linhas, colunas) ===")
print(df_imdb.shape)

# 1.4 Estatísticas básicas das colunas numéricas
print("\\n=== ESTATÍSTICAS DAS COLUNAS NUMÉRICAS ===")
print(df_imdb.describe())

# 1.5 Contando valores nulos por coluna
print("\\n=== VALORES NULOS POR COLUNA ===")
print(df_imdb.isnull().sum())'''
    },
    {
        "titulo_slide": "Script: Remoção de Duplicatas",
        "subtitulo": "script__duplicatas.py — Exercício 3",
        "codigo": '''# EXERCÍCIO 3 — Duplicatas

import pandas as pd

caminho = r"filmes_imdb\\dados\\dados_excel\\filmes_imdb.xlsx"
df_imdb = pd.read_excel(caminho)

# 3.1 Verificando duplicatas pela combinação título + ano
print("=== LINHAS DUPLICADAS (mesmo título e ano) ===")
duplicatas = df_imdb[df_imdb.duplicated(
    subset=['titulo', 'ano_lancamento'],
    keep='first'
)]
print(duplicatas[['id_filme', 'titulo', 'ano_lancamento',
                   'rating_imdb', 'duracao_min']])
print(f"\\nTotal de duplicatas encontradas: {len(duplicatas)}")

# 3.2 Removendo as duplicatas
df_imdb.drop_duplicates(
    subset=['titulo', 'ano_lancamento'],
    keep='first',
    inplace=True
)

# Resetando o índice após remoção
df_imdb.reset_index(drop=True, inplace=True)

print(f"\\nLinhas após remoção de duplicatas: {len(df_imdb)}")'''
    },
    {
        "titulo_slide": "Script: Padronização de Texto",
        "subtitulo": "script__padronizao.py — Exercício 4",
        "codigo": '''# EXERCÍCIO 4 — Padronização de Texto

import pandas as pd

caminho = r"filmes_imdb\\dados\\dados_excel\\filmes_imdb.xlsx"
df_imdb = pd.read_excel(caminho)

# 4.1 Verificando os valores únicos de 'genero'
print("=== GÊNEROS ÚNICOS (antes da limpeza) ===")
print(sorted(df_imdb['genero'].dropna().unique()))

# 4.2 Padronizando a coluna 'genero'
df_imdb['genero'] = df_imdb['genero'].str.strip().str.title()

# 4.3 Corrigindo abreviações e variações conhecidas
mapa_generos = {
    'Anim'            : 'Animation',
    'Scifi'           : 'Sci-Fi',
    'Science Fiction' : 'Sci-Fi',
    'Science-Fiction' : 'Sci-Fi',
    'Thriler'         : 'Thriller',
    'Acton'           : 'Action',
    'Fantasia'        : 'Fantasy',
}
df_imdb['genero'] = df_imdb['genero'].replace(mapa_generos)

print("\\n=== GÊNEROS ÚNICOS (após limpeza) ===")
print(sorted(df_imdb['genero'].dropna().unique()))

# 4.4 Padronizando 'nome_pais'
print("\\n=== PAÍSES ÚNICOS (antes) ===")
print(df_imdb['nome_pais'].dropna().unique())
df_imdb['nome_pais'] = df_imdb['nome_pais'].str.strip().str.title()
print("\\n=== PAÍSES ÚNICOS (após) ===")
print(df_imdb['nome_pais'].dropna().unique())'''
    },
    {
        "titulo_slide": "Script: Outliers e Valores Inválidos",
        "subtitulo": "script__outliers.py — Exercício 5",
        "codigo": '''# EXERCÍCIO 5 — Outliers e Valores Inválidos

import pandas as pd
import numpy as np

caminho = r"filmes_imdb\\dados\\dados_excel\\filmes_imdb.xlsx"
df_imdb = pd.read_excel(caminho)

# 5.1 Verificando o rating_imdb (range válido: 0 a 10)
print("=== RATINGS FORA DO RANGE (0 a 10) ===")
ratings_invalidos = df_imdb[
    (df_imdb['rating_imdb'] < 0) | (df_imdb['rating_imdb'] > 10)
]
print(ratings_invalidos[['id_filme', 'titulo', 'rating_imdb']])

df_imdb['rating_imdb'] = np.where(
    (df_imdb['rating_imdb'] < 0) | (df_imdb['rating_imdb'] > 10),
    np.nan, df_imdb['rating_imdb']
)
mediana_rating = df_imdb['rating_imdb'].median()
df_imdb['rating_imdb'] = df_imdb['rating_imdb'].fillna(mediana_rating)

# 5.2 Verificando 'duracao_min'
df_imdb['duracao_min'] = pd.to_numeric(
    df_imdb['duracao_min'], errors='coerce')
df_imdb['duracao_min'] = df_imdb['duracao_min'].replace(0, np.nan)
mediana_duracao = df_imdb['duracao_min'].median()
df_imdb['duracao_min'] = df_imdb['duracao_min'].fillna(mediana_duracao)

# 5.3 Verificando 'ano_lancamento'
df_imdb['ano_lancamento'] = pd.to_numeric(
    df_imdb['ano_lancamento'], errors='coerce')
df_imdb['ano_lancamento'] = np.where(
    (df_imdb['ano_lancamento'] < 1888) |
    (df_imdb['ano_lancamento'] > 2026),
    np.nan, df_imdb['ano_lancamento']
)'''
    },
    {
        "titulo_slide": "Script: Modificando Células com .loc[]",
        "subtitulo": "script_modificando_na.py — Métodos 1, 2 e 3",
        "codigo": '''# MODIFICANDO UMA CÉLULA ESPECÍFICA COM .loc[]
# Sintaxe: df.loc[índice_da_linha, 'nome_da_coluna'] = novo_valor

import pandas as pd
caminho = r"filmes_imdb\\dados\\dados_excel\\filmes_imdb.xlsx"
df_imdb = pd.read_excel(caminho)

# MÉTODO 1 — Usando o índice diretamente
df_imdb.loc[36, 'bilheteria_dolares'] = 321500000  # Inglourious Basterds
df_imdb.loc[42, 'duracao_min'] = 154               # Pulp Fiction
df_imdb.loc[46, 'data_lancamento'] = pd.Timestamp('1994-07-06')  # Forrest Gump
df_imdb.loc[49, 'titulo'] = 'Desconhecido'         # Linha sem título

# MÉTODO 2 — Localizando a linha pelo conteúdo
indice = df_imdb[df_imdb['titulo'] == 'Inglourious Basterds'].index[0]
df_imdb.loc[indice, 'bilheteria_dolares'] = 321500000

# Editando várias colunas de uma vez (dicionário)
indice_unknown = df_imdb[df_imdb['titulo'] == 'Unknown Movie'].index[0]
df_imdb.loc[indice_unknown,
    ['nome_pais', 'continente', 'idioma_principal']] = [
    'França', 'Europa', 'Francês'
]

# MÉTODO 3 — Editando SOMENTE SE a célula estiver vazia
if pd.isna(df_imdb.loc[36, 'bilheteria_dolares']):
    df_imdb.loc[36, 'bilheteria_dolares'] = 321500000
    print("Célula preenchida com sucesso!")
else:
    print(f"Já tem valor: {df_imdb.loc[36, 'bilheteria_dolares']}")'''
    },
    {
        "titulo_slide": "Script: Exportando o Dataset Limpo",
        "subtitulo": "script__salvando_.py — Exercício 6",
        "codigo": '''# EXERCÍCIO 6 — Exportando o Dataset Limpo

import pandas as pd
import numpy as np

caminho = r"filmes_imdb\\dados\\dados_excel\\filmes_imdb.xlsx"
df_imdb = pd.read_excel(caminho)

# --- Aplique todos os tratamentos dos exercícios 2 a 5 aqui ---

# Salvando como CSV
# index=False      → não salva o índice como coluna extra
# encoding='utf-8-sig' → garante que acentos apareçam corretamente
# sep=';'          → separador padrão brasileiro
df_imdb.to_csv(
    r"filmes_imdb\\dados\\dados_csv\\tab_imdb_limpo.csv",
    index=False,
    encoding='utf-8-sig',
    sep=';'
)

# Salvando como Excel
# sheet_name → define o nome da aba dentro do Excel
df_imdb.to_excel(
    r"filmes_imdb\\dados\\dados_excel\\tab_imdb_limpo.xlsx",
    index=False,
    sheet_name='filmes_limpo'
)

# Relatório final de conferência
print("=== RELATÓRIO FINAL DO DATASET LIMPO ===")
print(f"Total de linhas   : {df_imdb.shape[0]}")
print(f"Total de colunas  : {df_imdb.shape[1]}")
print(f"\\nValores nulos restantes:")
print(df_imdb.isnull().sum())
print("\\nArquivos salvos com sucesso!")'''
    },
]

# ============================================================
# CARREGAR APRESENTAÇÃO ORIGINAL
# ============================================================

prs = Presentation("R1S06.pptx")

# Capturar dimensões originais
slide_width = prs.slide_width
slide_height = prs.slide_height

# ============================================================
# FUNÇÃO: Criar slide de código respeitando layout original
# ============================================================

def add_code_slide(prs, titulo, subtitulo, codigo):
    """
    Adiciona um novo slide com caixa de título e bloco de código.
    Usa o layout em branco (último layout disponível) para não
    interferir com o design original.
    """
    # Tenta usar o layout 'Blank' ou o último disponível
    slide_layout = prs.slide_layouts[6]  # layout em branco padrão
    slide = prs.slides.add_slide(slide_layout)

    # ---------- FUNDO ESCURO (azul-escuro do tema original) ----------
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x06, 0x3E, 0x5F)  # #063E5F

    # ---------- BARRA DE TÍTULO (topo) ----------
    left   = Inches(0)
    top    = Inches(0)
    width  = slide_width
    height = Inches(1.3)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_tf  = title_box.text_frame
    title_tf.word_wrap = True

    # Linha 1: título principal
    p1 = title_tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    run1 = p1.add_run()
    run1.text = titulo
    run1.font.size   = Pt(22)
    run1.font.bold   = True
    run1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run1.font.name   = "Calibri"

    # Linha 2: subtítulo/nome do arquivo
    p2 = title_tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = subtitulo
    run2.font.size   = Pt(13)
    run2.font.bold   = False
    run2.font.color.rgb = RGBColor(0x42, 0x96, 0xCD)  # azul-claro do tema
    run2.font.name   = "Calibri"

    # ---------- BLOCO DE CÓDIGO ----------
    code_left   = Inches(0.3)
    code_top    = Inches(1.4)
    code_width  = slide_width  - Inches(0.6)
    code_height = slide_height - Inches(1.6)

    code_box = slide.shapes.add_textbox(code_left, code_top,
                                        code_width, code_height)
    # Fundo do bloco de código (cinza escuro)
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)

    tf = code_box.text_frame
    tf.word_wrap = False

    linhas = codigo.split("\n")
    for i, linha in enumerate(linhas):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = linha

        # Colorir comentários de forma diferente
        if linha.strip().startswith("#"):
            run.font.color.rgb = RGBColor(0x6A, 0x99, 0x55)  # verde
        elif linha.strip().startswith('"""') or linha.strip().startswith("'''"):
            run.font.color.rgb = RGBColor(0xCE, 0x91, 0x78)  # laranja
        else:
            run.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)  # branco-acinzentado

        run.font.size = Pt(9)
        run.font.name = "Courier New"
        run.font.bold = False

        # Espaçamento entre linhas
        p.space_before = Pt(0)
        p.space_after  = Pt(0)

    return slide


# ============================================================
# ADICIONAR SLIDES DE CÓDIGO À APRESENTAÇÃO
# ============================================================

print("Adicionando slides de código à apresentação...")

for i, script in enumerate(scripts_data):
    add_code_slide(
        prs,
        titulo    = script["titulo_slide"],
        subtitulo = script["subtitulo"],
        codigo    = script["codigo"]
    )
    print(f"  ✓ Slide adicionado: {script['titulo_slide']}")

# ============================================================
# SALVAR APRESENTAÇÃO
# ============================================================

output_path = "R1S06_com_scripts.pptx"
prs.save(output_path)

print(f"\n✅ Apresentação salva com sucesso: {output_path}")
print(f"   Total de slides: {len(prs.slides)}")
print(f"   (20 originais + {len(scripts_data)} novos slides de código)")